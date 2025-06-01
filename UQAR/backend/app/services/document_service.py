import os
import uuid
import logging
# Import chromadb conditionally to avoid errors
try:
    import chromadb
    CHROMADB_AVAILABLE = True
except (ImportError, RuntimeError) as e:
    logging.warning(f"ChromaDB not available: {e}")
    CHROMADB_AVAILABLE = False
from typing import List, Optional, Dict, Any
from fastapi import UploadFile
from sqlalchemy.orm import Session

from ..models.document import Document, DocumentStatus, DocumentType
from ..models.section import Section
from ..core.config import settings

logger = logging.getLogger(__name__)

class DocumentService:
    """Service pour la gestion des documents et leur vectorisation"""
    
    def __init__(self, db: Session):
        self.db = db
        self.upload_dir = settings.UPLOAD_DIR
        self.allowed_extensions = settings.get_allowed_extensions()
        self.max_file_size = settings.MAX_FILE_SIZE
        
        # Créer le répertoire d'upload s'il n'existe pas
        os.makedirs(self.upload_dir, exist_ok=True)
        
        # Only initialize ChromaDB if it's available
        self.chroma_client = None
        if CHROMADB_AVAILABLE:
            try:
                # Tenter d'initialiser ChromaDB HTTP
                self.chroma_client = chromadb.HttpClient(
                    host=settings.CHROMA_HOST,
                    port=settings.CHROMA_PORT
                )
                logger.info("ChromaDB initialized as HttpClient")
            except Exception as e:
                logger.warning(f"Failed to initialize ChromaDB HttpClient: {e}")
                # Fallback to PersistentClient
                self.chroma_client = chromadb.PersistentClient(
                    path=settings.CHROMA_PERSIST_DIRECTORY
                )
                logger.info("ChromaDB initialized as PersistentClient")
        else:
            logger.warning("ChromaDB is not available, document indexing features will be disabled")
    
    async def upload_document(self, file: UploadFile, section_id: int, user_id: int) -> Document:
        """
        Upload un document et crée l'entrée dans la base de données
        """
        # Vérifier si la section existe
        section = self.db.query(Section).filter(Section.id == section_id).first()
        if not section:
            raise ValueError(f"Section {section_id} non trouvée")
        
        # Vérifier si l'utilisateur est l'enseignant de cette section
        if section.teacher_id != user_id:
            raise ValueError("Vous n'êtes pas l'enseignant de cette section")
        
        # Vérifier l'extension du fichier
        file_ext = os.path.splitext(file.filename)[1].lower()
        if file_ext not in self.allowed_extensions:
            raise ValueError(f"Type de fichier non autorisé. Extensions acceptées: {', '.join(self.allowed_extensions)}")
        
        # Déterminer le type de document
        document_type = self._get_document_type(file_ext)
        
        # Générer un nom de fichier unique
        unique_filename = f"{uuid.uuid4().hex}{file_ext}"
        file_path = os.path.join(self.upload_dir, unique_filename)
        
        # Sauvegarder le fichier
        with open(file_path, "wb") as f:
            content = await file.read()
            
            # Vérifier la taille du fichier
            if len(content) > self.max_file_size:
                raise ValueError(f"Fichier trop volumineux. Maximum: {self.max_file_size / (1024 * 1024)}MB")
                
            f.write(content)
        
        # Créer l'entrée dans la base de données
        new_document = Document(
            filename=unique_filename,
            original_filename=file.filename,
            file_path=file_path,
            file_size=len(content),
            document_type=document_type,
            mime_type=file.content_type or "application/octet-stream",
            section_id=section_id,
            status=DocumentStatus.UPLOADED
        )
        
        self.db.add(new_document)
        self.db.commit()
        self.db.refresh(new_document)
        
        # Lancer le traitement du document en arrière-plan
        # Note: Dans une implémentation réelle, cela devrait être fait de manière asynchrone
        # avec une file d'attente comme Celery
        try:
            await self.process_document(new_document.id)
        except Exception as e:
            logger.error(f"Erreur lors du traitement du document {new_document.id}: {e}")
            # L'erreur est capturée mais on ne bloque pas l'upload
        
        return new_document
    
    async def process_document(self, document_id: int) -> Document:
        """
        Traite un document: extraction de texte et vectorisation
        """
        document = self.db.query(Document).filter(Document.id == document_id).first()
        if not document:
            raise ValueError(f"Document {document_id} non trouvé")
        
        # Mettre à jour le statut
        document.status = DocumentStatus.PROCESSING
        self.db.commit()
        
        try:
            # Extraire le texte selon le type de document
            extracted_text = self._extract_text(document)
            
            # Mettre à jour le document avec le texte extrait
            document.extracted_text = extracted_text
            document.text_length = len(extracted_text)
            
            # Compter le nombre de pages pour les PDF
            if document.document_type == DocumentType.PDF:
                document.page_count = self._count_pdf_pages(document.file_path)
            
            # Vectoriser le document dans ChromaDB
            chunks = self._chunk_text(extracted_text)
            self._vectorize_chunks(document, chunks)
            
            # Mettre à jour le statut
            document.status = DocumentStatus.PROCESSED
            document.is_vectorized = True
            document.vector_count = len(chunks)
            document.processed_at = func.now()
            
            self.db.commit()
            
            return document
            
        except Exception as e:
            # En cas d'erreur, mettre à jour le statut
            document.status = DocumentStatus.ERROR
            document.processing_error = str(e)
            self.db.commit()
            
            logger.error(f"Erreur lors du traitement du document {document_id}: {e}")
            raise
    
    def get_document(self, document_id: int) -> Optional[Document]:
        """
        Récupère un document par son ID
        """
        return self.db.query(Document).filter(Document.id == document_id).first()
    
    def get_section_documents(self, section_id: int) -> List[Document]:
        """
        Récupère tous les documents d'une section
        """
        return self.db.query(Document).filter(Document.section_id == section_id).all()
    
    def delete_document(self, document_id: int, user_id: int) -> bool:
        """
        Supprime un document et ses vecteurs
        """
        document = self.db.query(Document).filter(Document.id == document_id).first()
        if not document:
            return False
        
        # Vérifier si l'utilisateur est l'enseignant de cette section
        section = self.db.query(Section).filter(Section.id == document.section_id).first()
        if section.teacher_id != user_id:
            raise ValueError("Vous n'êtes pas l'enseignant de cette section")
        
        # Supprimer le fichier
        if os.path.exists(document.file_path):
            try:
                os.remove(document.file_path)
                logger.info(f"Fichier supprimé: {document.file_path}")
            except Exception as e:
                logger.error(f"Erreur lors de la suppression du fichier {document.file_path}: {e}")
        
        # Supprimer les vecteurs de ChromaDB
        if document.is_vectorized and self.chroma_client is not None:
            try:
                collection = self._get_chroma_collection(section.chroma_collection_name)
                if collection:
                    # Supprimer les vecteurs avec l'ID du document
                    collection.delete(filter={"document_id": str(document.id)})
                    logger.info(f"Vecteurs supprimés pour le document {document_id}")
            except Exception as e:
                logger.error(f"Erreur lors de la suppression des vecteurs du document {document_id}: {e}")
                # Ne pas bloquer la suppression en cas d'erreur avec ChromaDB
        
        # Supprimer le document de la base de données
        self.db.delete(document)
        self.db.commit()
        logger.info(f"Document {document_id} supprimé de la base de données")
        
        return True
    
    def _get_document_type(self, file_ext: str) -> DocumentType:
        """
        Détermine le type de document à partir de l'extension
        """
        ext_mapping = {
            ".pdf": DocumentType.PDF,
            ".docx": DocumentType.DOCX,
            ".pptx": DocumentType.PPTX,
            ".txt": DocumentType.TXT,
            ".md": DocumentType.MD
        }
        
        return ext_mapping.get(file_ext, DocumentType.TXT)
    
    def _extract_text(self, document: Document) -> str:
        """
        Extrait le texte d'un document selon son type
        """
        file_path = document.file_path
        
        if document.document_type == DocumentType.PDF:
            return self._extract_text_from_pdf(file_path)
        elif document.document_type == DocumentType.DOCX:
            return self._extract_text_from_docx(file_path)
        elif document.document_type == DocumentType.PPTX:
            return self._extract_text_from_pptx(file_path)
        elif document.document_type == DocumentType.TXT:
            return self._extract_text_from_txt(file_path)
        elif document.document_type == DocumentType.MD:
            return self._extract_text_from_txt(file_path)
        else:
            raise ValueError(f"Type de document non supporté: {document.document_type}")
    
    def _extract_text_from_pdf(self, file_path: str) -> str:
        """
        Extrait le texte d'un fichier PDF
        """
        from PyPDF2 import PdfReader
        
        reader = PdfReader(file_path)
        text = ""
        
        for page in reader.pages:
            text += page.extract_text() + "\n\n"
        
        return text
    
    def _extract_text_from_docx(self, file_path: str) -> str:
        """
        Extrait le texte d'un fichier DOCX
        """
        import docx
        
        doc = docx.Document(file_path)
        text = ""
        
        for para in doc.paragraphs:
            text += para.text + "\n"
        
        return text
    
    def _extract_text_from_pptx(self, file_path: str) -> str:
        """
        Extrait le texte d'un fichier PPTX
        """
        from pptx import Presentation
        
        pres = Presentation(file_path)
        text = ""
        
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            text += "\n"
        
        return text
    
    def _extract_text_from_txt(self, file_path: str) -> str:
        """
        Extrait le texte d'un fichier texte
        """
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    
    def _count_pdf_pages(self, file_path: str) -> int:
        """
        Compte le nombre de pages d'un PDF
        """
        from PyPDF2 import PdfReader
        
        reader = PdfReader(file_path)
        return len(reader.pages)
    
    def _chunk_text(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """
        Découpe le texte en chunks pour la vectorisation
        """
        chunks = []
        
        if len(text) <= chunk_size:
            chunks.append(text)
        else:
            for i in range(0, len(text), chunk_size - overlap):
                chunk = text[i:i + chunk_size]
                if chunk:
                    chunks.append(chunk)
        
        return chunks
    
    def _vectorize_chunks(self, document: Document, chunks: List[str]) -> None:
        """
        Vectorise les chunks dans ChromaDB
        """
        if not chunks:
            logger.warning(f"Aucun chunk à vectoriser pour le document {document.id}")
            return
            
        if self.chroma_client is None:
            logger.warning(f"ChromaDB client non initialisé, vectorisation impossible pour le document {document.id}")
            return
        
        try:
            # Récupérer la section pour obtenir le nom de la collection
            section = self.db.query(Section).filter(Section.id == document.section_id).first()
            if not section:
                logger.error(f"Section {document.section_id} non trouvée pour le document {document.id}")
                return
        
            # Récupérer ou créer la collection
            collection = self._get_chroma_collection(section.chroma_collection_name)
            
            # Vérifier si la collection est None (ChromaDB non disponible)
            if collection is None:
                logger.warning(f"ChromaDB collection {section.chroma_collection_name} non disponible, vectorisation ignorée")
                return
        
            # Préparer les données pour ChromaDB
            ids = [f"{document.id}_{i}" for i in range(len(chunks))]
            metadatas = [{
                "document_id": str(document.id),
                "filename": document.original_filename,
                "chunk_index": i,
                "document_type": document.document_type.value,
                "section_id": str(document.section_id)
            } for i in range(len(chunks))]
            
            # Vérifier si des chunks avec les mêmes IDs existent déjà
            try:
                existing_ids = collection.get(ids=ids, include=["documents"])
                if existing_ids and "ids" in existing_ids and existing_ids["ids"]:
                    logger.info(f"Suppression des chunks existants pour le document {document.id}")
                    collection.delete(ids=ids)
            except Exception as e:
                logger.warning(f"Impossible de vérifier les chunks existants: {e}, continuons avec l'ajout")
            
            # Ajouter les chunks à ChromaDB, en lots si nécessaire
            batch_size = 100  # Réduire la taille des lots si nécessaire
            for i in range(0, len(chunks), batch_size):
                batch_ids = ids[i:i+batch_size]
                batch_chunks = chunks[i:i+batch_size]
                batch_metadatas = metadatas[i:i+batch_size]
                
                try:
                    collection.add(
                        ids=batch_ids,
                        documents=batch_chunks,
                        metadatas=batch_metadatas
                    )
                    logger.info(f"Vectorisation réussie pour le lot {i//batch_size + 1}/{(len(chunks) + batch_size - 1)//batch_size} du document {document.id}")
                except Exception as batch_error:
                    logger.error(f"Erreur lors de la vectorisation du lot {i//batch_size + 1}: {batch_error}")
            
            logger.info(f"Vectorisation complétée pour le document {document.id} ({len(chunks)} chunks)")
        except Exception as e:
            logger.error(f"Erreur lors de la vectorisation dans ChromaDB: {e}")
            # Ne pas bloquer le traitement en cas d'erreur avec ChromaDB
    
    def _get_chroma_collection(self, collection_name: str):
        """
        Récupère ou crée une collection ChromaDB
        """
        if self.chroma_client is None:
            logger.warning(f"ChromaDB client not initialized, cannot get collection {collection_name}")
            return None
            
        try:
            # Vérifier si la collection existe
            try:
                collection = self.chroma_client.get_collection(name=collection_name)
                logger.info(f"Collection ChromaDB {collection_name} récupérée avec succès")
                return collection
            except Exception as e:
                logger.warning(f"Collection {collection_name} non trouvée, tentative de création: {e}")
            # Tenter de créer la collection si elle n'existe pas
            collection = self.chroma_client.create_collection(
                name=collection_name,
                metadata={"hnsw:space": "cosine"}
            )
            logger.info(f"Collection ChromaDB {collection_name} créée avec succès")
            return collection
        except Exception as e:
            logger.error(f"Impossible d'accéder ou de créer la collection {collection_name}: {e}")
            return None

# Import nécessaire pour le timestamp
from sqlalchemy.sql import func 